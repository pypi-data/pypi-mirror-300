from dataclasses import dataclass
from pptx import Presentation
from google.cloud import texttospeech_v1beta1 as tts
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
import os
import re

@dataclass
class Meta:
    # PPT settings
    ppt_file: str 
    ppt_path: str = 'data/ppt/'  # Directory for the PPT and image files
    image_prefix: str = 'slide'  # The prefix for image file names (used when saving slides as images)
    image_extension: str = 'PNG'  # The image file format (default is PNG)
    ppt_extension: str = '.pptx'  # The PowerPoint file extension

    # Google TTS settings
    voice_enabled: bool = True  # Enable or disable voice narration
    google_application_credentials: str = None  # Location of the Google API key (downloaded as JSON)
    voice_path: str = 'data/voice/'  # Directory to save the generated audio files
    max_size: int = 4500  # Maximum text size limit for a single Google TTS API request (default 5000)
    slide_break: float = 2  # Time delay (in seconds) between slides
    line_break: float = 0.7  # Time delay (in seconds) when there's a line break in the text (e.g., '\n')
    lang: str = 'E'  # Language setting: 'E' for English, 'K' for Korean 
    wave: bool = False  # Whether to use Wavenet voices (True or False)

    # MoviePy video settings
    fps: int = 24  # Frames per second for the video


def ppt_to_video(meta: Meta): 
    if not os.path.exists(meta.ppt_path):
        os.makedirs(meta.ppt_path)

    if meta.voice_enabled:
        if meta.google_application_credentials == None:
            print('*****')
            print('Need to set up Google Cloud Authentication')
            print('Please refer to the README.md')
            print('*****')
            return None

        if not os.path.exists(meta.voice_path):
            os.makedirs(meta.voice_path)
        num = ppt_to_text(meta)
        timepoints = ppt_tts(meta, num)
        video_from_ppt_and_voice(meta, timepoints)
    else:
        num = ppt_to_text(meta)
        video_from_ppt(meta, num)

def _clean_text(input_text):
    # Ensure UTF-8 compatibility: decode and encode to handle encoding correctly
    input_text = input_text.encode('utf-8').decode('utf-8')
    
    # 1. Replace multiple spaces with a single space
    input_text = re.sub(r'\s+', ' ', input_text)
    
    # 2. Remove non-Korean, non-English chars, non-numbers, and special characters except commas, periods, question marks, exclamation marks, and spaces
    input_text = re.sub(r'[^a-zA-Z0-9가-힣.,?!\n\s]', '', input_text)
    
    # 3. Replace multiple newlines with a single newline
    input_text = re.sub(r'(\n)+', '\n', input_text)
    
    # 4. Collapse spaces between \n and \n into a single \n
    input_text = re.sub(r'(?<=\n)\s+(?=\n)', '', input_text)
    
    # 5. Remove any trailing newline at the end of the text
    input_text = input_text.rstrip('\n')
    
    # 6. Remove any space before a newline
    input_text = re.sub(r'\s+(?=\n)', '', input_text)
    
    # Return cleaned text
    return input_text.strip()

def _write_to_file(content, current_file_number, current_size, meta: Meta):
    txt_file = f"{os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, ''))}_{current_file_number}.txt"

    mode = 'w' if current_size == 0 else 'a'
    with open(txt_file, mode, encoding='utf-8') as notes_file:
        notes_file.write(content)
    
    return current_size + len(content.encode('utf-8'))

def ppt_to_text(meta: Meta):
    ppt = Presentation(os.path.join(meta.ppt_path, meta.ppt_file))
    if not meta.voice_enabled:
        return len(ppt.slides)

    header = '''<speak>\n'''
    footer = '''</speak>'''
    file_number = 0
    current_size = _write_to_file(header, file_number, 0, meta)

    mark_separator = '.' # for the Google TTS English engine, MARK tag seems need to be followed by a char.
    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            notes = _clean_text(notes)
            slide_content = f'<mark name="slide{slide_number}"/>{mark_separator}\n<break time="{round(meta.slide_break/2,1)}s"/>\n'
            slide_content += notes.replace('\n', f'\n<break time="{meta.line_break}s"/>\n') + f'\n<break time="{meta.slide_break}s"/>\n'
        else:
            slide_content = f'<mark name="slide{slide_number}"/>{mark_separator}\n<break time="{meta.slide_break}s"/>\n'

        if current_size + len(slide_content.encode('utf-8')) > meta.max_size:
            _write_to_file(footer, file_number, current_size, meta)
            file_number += 1
            current_size = 0
            slide_content = header + slide_content

        current_size = _write_to_file(slide_content, file_number, current_size, meta)

    _write_to_file(footer, file_number, current_size, meta)
    txt_file_number = file_number+1

    return txt_file_number 


def ppt_tts(meta: Meta, txt_file_number: int):

    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = meta.google_application_credentials

    client = tts.TextToSpeechClient()
    language_code = 'ko-KR' if meta.lang == 'K' else 'en-US' 
    tag = 'D' if meta.lang == 'K' else 'B' 
    name = language_code+'-Wavenet-'+tag 
    if meta.wave == True: # WaveNet voice (1 mil words/month vs 4 mil in basic)
        voice = tts.VoiceSelectionParams(language_code=language_code, name=name, ssml_gender=tts.SsmlVoiceGender.MALE)
    else:
        voice = tts.VoiceSelectionParams(language_code=language_code, ssml_gender=tts.SsmlVoiceGender.MALE)
    audio_config = tts.AudioConfig(audio_encoding=tts.AudioEncoding.MP3)
    
    timepoint_dict = {}
    for i in range(txt_file_number):
        txt_file = f"{os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, '_'+str(i)+'.txt'))}"
        voice_file =os.path.join(meta.voice_path, meta.ppt_file.replace(meta.ppt_extension, '_'+str(i)+'.mp3'))

        with open(txt_file, 'r', encoding='utf-8') as file:
            text_content = file.read()

        synthesis_input = tts.SynthesisInput(ssml=text_content)
        request = tts.SynthesizeSpeechRequest(
            input=synthesis_input,
            voice=voice,
            audio_config=audio_config, 
            enable_time_pointing=[tts.SynthesizeSpeechRequest.TimepointType.SSML_MARK]
        )
        response = client.synthesize_speech(request=request)

        with open(voice_file, "wb") as out:
            out.write(response.audio_content)
            print(voice_file + ' done')

        timepoint_list = []
        if response.timepoints:
            for time_point in response.timepoints:
                print(f'Mark name: {time_point.mark_name}, Time: {time_point.time_seconds} seconds')
                timepoint_list.append([int(time_point.mark_name[5:]), time_point.time_seconds])
        else:
            print('No timepoints found.')
        timepoint_dict[voice_file] = timepoint_list

    return timepoint_dict

def video_from_ppt_and_voice(meta: Meta, timepoints, fps=24):
    images_path = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension,''))
    output_file = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '.mp4'))
    video_with_audios = []

    for audio_file, slide_times in timepoints.items():
        audio_clip = AudioFileClip(audio_file)

        video_clips = []
        for i in range(len(slide_times)):
            start_time = slide_times[i][1]  # Get the start time for the slide
            if i < len(slide_times)-1:
                end_time = slide_times[i + 1][1]  # Get the end time for the next slide
            else:
                end_time = audio_clip.duration
            slide_number = slide_times[i][0]

            # Construct the image filename
            slide_image_filename = f'{meta.image_prefix}{slide_number}.PNG'
            slide_image_path = os.path.join(images_path, slide_image_filename)

            # Load the slide image
            slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)
            video_clips.append(slide_clip)

        # Concatenate video clips for the current audio
        video_with_audio = concatenate_videoclips(video_clips)
        video_with_audio = video_with_audio.set_audio(audio_clip).volumex(2)
        video_with_audios.append(video_with_audio)

    # Concatenate all videos into one final video
    final_video = concatenate_videoclips(video_with_audios)

    # Set fps for the final video
    final_video.fps = fps
    
    # final_video.write_videofile(output_file, codec="libx264")
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video with audio generated and saved')

def video_from_ppt(meta: Meta, num_slides: int, fps=24):
    images_path = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension,''))
    output_file = os.path.join(meta.ppt_path, meta.ppt_file.replace(meta.ppt_extension, '.mp4'))

    video_clips = []
    for i in range(num_slides):
        start_time = i*meta.slide_break
        end_time = start_time+meta.slide_break

        slide_image_filename = f'{meta.image_prefix}{i}.{meta.image_extension}'
        slide_image_path = os.path.join(images_path, slide_image_filename)

        slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)
        video_clips.append(slide_clip)

    final_video = concatenate_videoclips(video_clips)

    # Set fps for the final video
    final_video.fps = fps
    
    # final_video.write_videofile(output_file, codec="libx264")
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video generated and saved')
